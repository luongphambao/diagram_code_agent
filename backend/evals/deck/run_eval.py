"""Deck eval runner — deterministic, no agent/LLM, CI-safe.

    cd backend
    uv run python -m evals.deck.run_eval
    uv run python -m evals.deck.run_eval --gate
    uv run python -m evals.deck.run_eval --update-baseline
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

_BACKEND = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(_BACKEND / "src"))
sys.path.insert(0, str(_BACKEND))

from memory.stores.csm import Evidence  # noqa: E402
from memory.stores.csm_adapter import from_artifacts  # noqa: E402
from domain.deck.deck import DeckPlan, build_deck_plan, score_deck_structure, validate_deck  # noqa: E402
from evals._core import (  # noqa: E402
    compare_to_baseline,
    load_cases,
    print_table,
    run_all_sync,
    write_baseline,
    write_results,
)
from domain.deck.deck_visual_qa import audit_pptx_deterministic  # noqa: E402
from evals.deck.judge import (  # noqa: E402
    METRIC_KEYS,
    STRUCTURE_METRIC_KEYS,
    VISUAL_AUDIT_METRIC_KEYS,
    score_deck,
    score_structure,
    score_visual_audit,
)

ALL_METRIC_KEYS = METRIC_KEYS + STRUCTURE_METRIC_KEYS + VISUAL_AUDIT_METRIC_KEYS

_DIR = Path(__file__).resolve().parent
_DATASET = _DIR / "dataset"
_RESULTS = _DIR / "results.json"
_BASELINE = _DIR / "baseline.json"

_COLUMNS = [
    ("case_id", "case"),
    ("scores.finding_recall", "recall"),
    ("scores.match", "match"),
    ("scores.n_findings", "n_find"),
    ("scores.score_ok", "struct_ok"),
    ("scores.kw_recall", "kw_rec"),
    ("scores.structural_score", "struct_score"),
    ("scores.issue_recall", "vis_recall"),
    ("scores.passed_ok", "vis_passed"),
]


def _build_plan(case: dict) -> tuple[DeckPlan, object]:
    """Build the model + storyboard for a case, applying any seeded-defect directives."""
    brief = case.get("brief", {})
    blueprint = case.get("blueprint", {})
    wbs = case.get("wbs", {})
    model = from_artifacts(brief, blueprint, wbs, tech_stack=case.get("tech_stack", {}))

    # Inject grounded evidence (evidence lives in evidence_log, not the artifacts).
    for i, ev in enumerate(case.get("evidence", []), start=1):
        model.evidence.append(Evidence(id=ev.get("id", f"EVD-{i}"),
                                       claim=ev.get("claim", ""),
                                       source_url=ev.get("source_url", "")))

    plan = build_deck_plan(model, wbs=wbs, brief=brief,
                           has_diagram=case.get("has_diagram", True))

    # Seeded-defect directives (keep golden cases tiny).
    if case.get("append_ghost_ref"):
        for s in plan.slides:
            if s.source_refs:
                s.source_refs.append("COMP-ghost-does-not-exist")
                break
    if case.get("drop_required_roles"):
        drop = set(case["drop_required_roles"])
        plan.slides = [s for s in plan.slides if s.narrative_role not in drop]
    return plan, model


def _apply_structure_defects(plan: DeckPlan, case: dict) -> None:
    """Mutate the plan in place to seed structure defects declared in the case."""
    if case.get("make_long_title"):
        for s in plan.slides:
            if s.title and "Head Page" not in s.layout:
                s.title = "A" * 100  # overlong title
                break
    if case.get("make_text_dense"):
        for s in plan.slides:
            if s.bullets and "Head Page" not in s.layout:
                s.bullets = [f"bullet {i}" for i in range(12)]  # too many bullets
                break
    if case.get("make_ungrounded_client_facing"):
        for s in plan.slides:
            if s.client_facing:
                s.source_refs = []  # clear source_refs on client-facing slide
                break


def _run_one(case: dict) -> dict:
    # Visual-audit cases (case type "visual_audit") are self-contained — they
    # build a minimal PPTX in memory and run audit_pptx_deterministic directly.
    if case.get("case_type") == "visual_audit":
        return _run_visual_case(case)

    plan, model = _build_plan(case)
    _apply_structure_defects(plan, case)
    findings = validate_deck(plan, model)
    struct = score_deck_structure(plan)
    scores = score_deck(findings, case)
    scores.update(score_structure(struct, case))
    # Visual audit metrics default to 1.0 for non-visual cases (not applicable)
    scores.setdefault("issue_recall", 1.0)
    scores.setdefault("passed_ok", 1.0)
    return {"case_id": case.get("id", "unknown"), "scores": scores}


def _run_visual_case(case: dict) -> dict:
    """Build a minimal python-pptx presentation per the case spec and audit it.

    Uses the "Title and Content" layout (index 1 of the default template) so
    that proper placeholder shapes (idx 0 = title, idx 1 = body) are present
    and can be detected by audit_pptx_deterministic.
    """
    import tempfile
    from pathlib import Path as _Path

    from pptx import Presentation  # type: ignore
    from pptx.util import Pt  # type: ignore

    defects = case.get("defects", [])
    prs = Presentation()
    # Layout 1 = "Title and Content" in the default template — has title (idx 0)
    # and body (idx 1) placeholders.
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # Fill title placeholder
    title_ph = slide.placeholders[0]
    title_ph.text = "A" * 90 if "long_title" in defects else case.get("slide_title", "Test Slide")

    # Fill body placeholder with many bullets if requested
    if "dense_bullets" in defects and len(slide.placeholders) > 1:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.text = "Bullet 0"
        for i in range(1, 12):
            p = tf.add_paragraph()
            p.text = f"Bullet {i}"

    # Add a tiny-font run if requested
    if "tiny_font" in defects and len(slide.placeholders) > 1:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        if not tf.text.strip():
            tf.text = "tiny text"
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.font.size = Pt(5)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name

    try:
        result = audit_pptx_deterministic(tmp_path)
        scores = score_visual_audit(result.model_dump(), case)
        # Not applicable for visual-only cases — default to passing
        scores.setdefault("finding_recall", 1.0)
        scores.setdefault("match", 1.0)
        scores.setdefault("n_findings", 0)
        scores.setdefault("score_ok", 1.0)
        scores.setdefault("kw_recall", 1.0)
        scores.setdefault("structural_score", 100)
    finally:
        _Path(tmp_path).unlink(missing_ok=True)

    return {"case_id": case.get("id", "unknown"), "scores": scores}


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", default=None)
    ap.add_argument("--gate", action="store_true")
    ap.add_argument("--update-baseline", action="store_true")
    args = ap.parse_args()

    cases = load_cases(_DATASET, args.case)
    if not cases:
        print(f"No deck cases in {_DATASET}")
        sys.exit(1)

    results = run_all_sync(cases, _run_one)
    print_table(results, _COLUMNS)
    write_results(_RESULTS, results)

    if args.update_baseline:
        write_baseline(_BASELINE, results, ALL_METRIC_KEYS)
        print(f"Baseline updated: {_BASELINE}")
        return

    passed, regressions = compare_to_baseline(results, _BASELINE, ALL_METRIC_KEYS)
    if regressions:
        print("REGRESSIONS:")
        for r in regressions:
            print(f"  {r['metric']}: {r['baseline']} -> {r['current']} (drop {r['drop']})")
    else:
        print("No regression vs baseline." if _BASELINE.exists() else "No baseline yet.")
    if args.gate and not passed:
        sys.exit(1)


if __name__ == "__main__":
    main()
