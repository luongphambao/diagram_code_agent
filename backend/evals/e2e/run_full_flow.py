"""Live end-to-end smoke test: drives the REAL agent through the full
deliverable pipeline — diagram -> PDF -> WBS -> PPT -> email -> Google
Calendar/Meet — auto-approving every HITL gate, and verifies each stage
produced what it promises.

This is intentionally NOT part of ``evals.run_all`` (the deterministic CI
gate). It calls a real LLM (OpenAI, via config.yaml) and, since Composio is
configured to run live, it SENDS A REAL EMAIL and CREATES A REAL GOOGLE
CALENDAR EVENT/MEET every time it runs. Run it explicitly when you want a
full-pipeline confidence check:

    cd backend
    uv run python -m evals.e2e.run_full_flow
    uv run python -m evals.e2e.run_full_flow --recipient someone@example.com

Exit code 0 = every stage passed (usable as a CI/CD gate). Exit code 1 =
at least one stage failed; the printed summary table says which.
"""

from __future__ import annotations

import argparse
import asyncio
import os
import re
import shutil
import sys
from pathlib import Path
from types import SimpleNamespace

# Allow running from the backend/ directory (mirrors evals/diagram/target.py).
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "src"))

from dotenv import load_dotenv

load_dotenv()
load_dotenv(Path(__file__).resolve().parents[2] / ".env")

from langchain_core.messages import HumanMessage, ToolMessage
from langgraph.types import Command

from agent import RECURSION_LIMIT, build_agent
from backends import reset_current_workspace, resolve_workspace, set_current_workspace
from context import SessionContext
from tools import GATE_TOOL_NAMES, clear_stage_markers

_REPO_ROOT = Path(__file__).resolve().parents[3]
_DOC_PATH = _REPO_ROOT / "example_doc" / "tom_tat_FMCG_Finance_Automation_Storybook.md"
_THREAD_ID = "e2e-fullflow"
_DEFAULT_RECIPIENT = "bao.luong@bnksolution.com"
_REQUIRED_ENV = [
    "OPENAI_API_KEY",
    "COMPOSIO_API_KEY",
    "GMAIL_CONNECTED_ACCOUNT_ID",
    "GOOGLE_CALENDAR_CONNECTED_ACCOUNT_ID",
    "GOOGLE_MEET_CONNECTED_ACCOUNT_ID",
]


# ---------------------------------------------------------------------------
# Gate/interrupt auto-approve loop (extends evals/diagram/target.py's pattern
# with SessionContext propagation and the slot_picker interrupt shape).
# ---------------------------------------------------------------------------

class Budget:
    """Shared, mutable resume counter — one budget for the whole run."""

    def __init__(self, n: int) -> None:
        self.remaining = n


async def _drain(agen) -> None:
    async for _ in agen:
        pass


async def _get_pending_interrupt(agent, config: dict):
    try:
        st = await agent.aget_state(config)
    except Exception:
        return None
    for task in getattr(st, "tasks", None) or []:
        for intr in getattr(task, "interrupts", None) or []:
            return getattr(intr, "value", intr)
    for intr in getattr(st, "interrupts", None) or []:
        return getattr(intr, "value", intr)
    return None


def _gate_name(val) -> str | None:
    if isinstance(val, dict):
        ars = val.get("action_requests") or []
        if ars:
            return ars[0].get("name")
    return None


def _resume_payload(val) -> dict:
    """Always-approve resume payload for either interrupt shape."""
    if isinstance(val, dict) and val.get("type") == "slot_picker":
        slots = val.get("slots") or []
        return {"decisions": [{"type": "approve", "selected_slot": slots[0] if slots else {}}]}
    return {"decisions": [{"type": "approve"}]}


async def drive_turn(agent, config: dict, ctx: SessionContext, text: str, budget: Budget) -> str:
    """Send one user turn, auto-approving every gate/interrupt it raises.

    Returns "" if the turn settled cleanly, else a short note explaining why
    it stopped early (budget exhausted / unrecognized interrupt).
    """
    await _drain(agent.astream(
        {"messages": [HumanMessage(content=text)]}, config,
        context=ctx, stream_mode=["updates"],
    ))
    while True:
        val = await _get_pending_interrupt(agent, config)
        if val is None:
            return ""
        is_slot_picker = isinstance(val, dict) and val.get("type") == "slot_picker"
        name = _gate_name(val)
        if not is_slot_picker and name not in GATE_TOOL_NAMES:
            return f"unrecognized interrupt (name={name!r}, keys={list(val) if isinstance(val, dict) else val!r})"
        if budget.remaining <= 0:
            return "resume budget exhausted"
        budget.remaining -= 1
        await _drain(agent.astream(
            Command(resume=_resume_payload(val)), config,
            context=ctx, stream_mode=["updates"],
        ))


def _tool_message_texts(state) -> list[str]:
    msgs = (getattr(state, "values", None) or {}).get("messages", [])
    texts = []
    for m in msgs:
        if not isinstance(m, ToolMessage):
            continue
        c = m.content
        if isinstance(c, list):
            c = " ".join(
                (blk.get("text", "") if isinstance(blk, dict) else str(blk)) for blk in c
            )
        texts.append(str(c))
    return texts


# ---------------------------------------------------------------------------
# Stage verification helpers
# ---------------------------------------------------------------------------

def _verify_file(ws: Path, name: str, *, magic: bytes | None = None) -> tuple[bool, str]:
    p = ws / name
    if not p.exists():
        return False, f"{name} was not created"
    data = p.read_bytes()
    if not data:
        return False, f"{name} exists but is empty"
    if magic and not data.startswith(magic):
        return False, f"{name} exists but does not start with expected header ({data[:8]!r})"
    return True, f"{name} OK ({len(data)} bytes)"


def _verify_pptx(ws: Path) -> tuple[bool, str]:
    ok, detail = _verify_file(ws, "out.pptx", magic=b"PK")
    if not ok:
        return ok, detail
    try:
        from pptx import Presentation
        prs = Presentation(str(ws / "out.pptx"))
        if len(prs.slides) < 2:
            return False, f"out.pptx has only {len(prs.slides)} slide(s)"
    except Exception as exc:  # noqa: BLE001
        return False, f"out.pptx exists but failed to open: {exc}"
    return True, f"{detail}, {len(prs.slides)} slides"


def _contains(texts: list[str], needle: str) -> bool:
    return any(needle in t for t in texts)


def _extract_first_record_name(text: str) -> str | None:
    m = re.search(r"-\s*(\S+)\s", text)
    return m.group(1) if m else None


# ---------------------------------------------------------------------------
# Main flow
# ---------------------------------------------------------------------------

async def _run(recipient: str) -> list[tuple[str, str, str]]:
    """Returns a list of (stage_name, status, detail); status in PASS/FAIL/SKIP."""
    results: list[tuple[str, str, str]] = []

    ws = resolve_workspace(_THREAD_ID)
    if ws.exists():
        shutil.rmtree(ws)
    ws.mkdir(parents=True, exist_ok=True)
    token = set_current_workspace(ws)
    try:
        clear_stage_markers()

        doc_text = _DOC_PATH.read_text(encoding="utf-8")
        (ws / "requirements.md").write_text(
            f"<untrusted_document>\n{doc_text}\n</untrusted_document>", encoding="utf-8",
        )

        agent = build_agent()
        config = {
            "configurable": {"thread_id": _THREAD_ID},
            "recursion_limit": RECURSION_LIMIT,
            "run_name": "e2e-full-flow",
        }
        ctx = SessionContext(
            user_email=recipient,
            composio_api_key=os.environ.get("COMPOSIO_API_KEY", ""),
            gmail_account_id=os.environ.get("GMAIL_CONNECTED_ACCOUNT_ID", ""),
            calendar_account_id=os.environ.get("GOOGLE_CALENDAR_CONNECTED_ACCOUNT_ID", ""),
            meet_account_id=os.environ.get("GOOGLE_MEET_CONNECTED_ACCOUNT_ID", ""),
        )
        budget = Budget(40)

        stages = [
            (
                "diagram",
                "The detailed requirements are saved in requirements.md in your working "
                "directory — read that file first, then design and render the architecture "
                "diagram for this FMCG finance automation proposal.",
                lambda texts: _verify_file(ws, "out.png"),
            ),
            (
                "pdf_report",
                "Generate the PDF report now — call generate_pdf_report with no arguments.",
                lambda texts: _verify_file(ws, "out.pdf", magic=b"%PDF"),
            ),
            (
                "wbs_excel",
                "Create the WBS and estimate effort for this project, then export the WBS Excel file.",
                lambda texts: _verify_file(ws, "wbs_filled.xlsx", magic=b"PK"),
            ),
            (
                "ppt_proposal",
                "Make the PPT proposal deck now.",
                lambda texts: _verify_pptx(ws),
            ),
            (
                "send_email",
                f"Send all the generated deliverables by email to {recipient}.",
                lambda texts: (
                    _contains(texts, "Email sent successfully"),
                    "found 'Email sent successfully' in tool output" if _contains(texts, "Email sent successfully")
                    else "no ToolMessage contained 'Email sent successfully'",
                ),
            ),
            (
                "schedule_meeting",
                "Propose meeting slots, then schedule a client meeting titled "
                f"'FMCG Automation Kickoff' with {recipient} using the first available slot.",
                lambda texts: (
                    (ws / "last_meeting.json").exists() and _contains(texts, "Meeting scheduled"),
                    "last_meeting.json written and 'Meeting scheduled' seen"
                    if (ws / "last_meeting.json").exists() and _contains(texts, "Meeting scheduled")
                    else f"last_meeting.json exists={(ws / 'last_meeting.json').exists()}, "
                         f"'Meeting scheduled' seen={_contains(texts, 'Meeting scheduled')}",
                ),
            ),
        ]

        aborted = False
        for name, message, verify in stages:
            if aborted:
                results.append((name, "SKIP", "earlier stage crashed the agent run"))
                continue
            try:
                note = await drive_turn(agent, config, ctx, message, budget)
                state = await agent.aget_state(config)
                texts = _tool_message_texts(state)
                ok, detail = verify(texts)
                if note and not ok:
                    detail = f"{detail} (turn did not settle cleanly: {note})"
                results.append((name, "PASS" if ok else "FAIL", detail))
            except Exception as exc:  # noqa: BLE001 — a crashed turn must not kill the whole report
                results.append((name, "FAIL", f"turn raised {type(exc).__name__}: {exc}"))
                aborted = True

        # --- Google Meet read-only lookup tools: call directly, not via the LLM ---
        import tools as tools_mod

        runtime = SimpleNamespace(context=ctx)
        try:
            records_text = tools_mod.list_meeting_records.func(runtime=runtime)
            records_ok = not records_text.startswith("ERROR")
            results.append((
                "meet_list_records", "PASS" if records_ok else "FAIL",
                records_text.splitlines()[0] if records_text else "(empty)",
            ))
        except Exception as exc:  # noqa: BLE001
            records_ok, records_text = False, ""
            results.append(("meet_list_records", "FAIL", f"raised {type(exc).__name__}: {exc}"))

        record_name = _extract_first_record_name(records_text) if records_ok else None
        if not record_name:
            results.append((
                "meet_transcript_recordings_participants", "SKIP",
                "no conference record yet (expected — nobody has joined the newly "
                "created meeting); list_meeting_records itself ran successfully",
            ))
        else:
            for stage_name, fn in (
                ("meet_transcript", tools_mod.get_meeting_transcript),
                ("meet_recordings", tools_mod.get_meeting_recordings),
                ("meet_participants", tools_mod.list_meeting_participants),
            ):
                try:
                    text = fn.func(conference_record_name=record_name, runtime=runtime)
                    ok = not text.startswith("ERROR")
                    results.append((stage_name, "PASS" if ok else "FAIL", text.splitlines()[0] if text else "(empty)"))
                except Exception as exc:  # noqa: BLE001
                    results.append((stage_name, "FAIL", f"raised {type(exc).__name__}: {exc}"))

        return results
    finally:
        reset_current_workspace(token)


def _preflight() -> list[str]:
    missing = [k for k in _REQUIRED_ENV if not os.environ.get(k)]
    if not _DOC_PATH.exists():
        missing.append(f"input doc missing: {_DOC_PATH}")
    return missing


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--recipient",
        default=os.environ.get("E2E_RECIPIENT_EMAIL", _DEFAULT_RECIPIENT),
        help="Email/attendee used for the send_email and schedule_meeting stages "
             f"(default: {_DEFAULT_RECIPIENT}).",
    )
    args = parser.parse_args()

    missing = _preflight()
    if missing:
        print("E2E full-flow ABORTED — missing prerequisites:")
        for m in missing:
            print(f"  - {m}")
        return 1

    print(f"Running full-pipeline E2E smoke test (recipient={args.recipient})...")
    print("This calls a real LLM and real Composio (Gmail + Google Calendar) — expect a few minutes.\n")

    results = asyncio.run(_run(args.recipient))

    width = max(len(name) for name, _, _ in results)
    print("\n=== E2E full-flow summary ===")
    for name, status, detail in results:
        print(f"  [{status:<4}] {name:<{width}}  {detail}")

    failed = [r for r in results if r[1] == "FAIL"]
    if failed:
        print(f"\n{len(failed)} stage(s) FAILED.")
        return 1
    print("\nAll stages PASSED.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
