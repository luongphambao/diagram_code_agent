"""Deck visual QA (Phase 3) — deterministic PPTX audit + targeted patcher.

No rendering dependency (no aspose, no LibreOffice): inspects the python-pptx
XML model directly.  Catches common layout issues before the user sees the file:

  * title_too_long   — title placeholder text > 80 chars (wraps badly on slides)
  * too_many_bullets — body has > 8 paragraph runs (text density / overflow risk)
  * table_overflow   — table cell text > 250 chars
  * tiny_font        — any run with an explicit font size < 8 pt (unreadable)
  * font_drift       — font family not in the BnK palette {Calibri, Calibri Light}
  * empty_body       — content slide with an empty body placeholder (no text / image)

`audit_pptx_deterministic` returns a `DeckVisualAuditResult` — structured,
serialisable, no side effects.

`patch_pptx_overflow` applies surgical text truncations and writes
`out_patched.pptx` next to the source file.  Used by the `generate_ppt_proposal`
gate tool when HIGH issues are found (docx §8.4 "targeted patch").
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Literal, Optional

from pydantic import BaseModel, Field

# python-pptx slide layout names that represent section dividers (no body expected).
_SECTION_LAYOUTS = frozenset(
    {"Cover-01", "Head Page", "Head-01", "BnK", "C2 -  Separator/ Dark", "Empty"}
)

_BNK_FONTS = frozenset({"Calibri", "Calibri Light"})

_MAX_TITLE_LEN = 80
_MAX_BULLETS = 8
_MAX_CELL_LEN = 250
_MIN_FONT_PT = 8.0

VISUAL_AUDIT_NAME = "deck_visual_audit.json"


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------

class SlideIssue(BaseModel):
    slide_idx: int
    slide_title: str
    issue_type: Literal[
        "title_too_long", "too_many_bullets", "table_overflow",
        "tiny_font", "font_drift", "empty_body",
    ]
    severity: Literal["high", "medium", "low"]
    detail: str
    patch_hint: Optional[str] = None


class DeckVisualAuditResult(BaseModel):
    pptx_path: str
    slide_count: int
    issues: list[SlideIssue] = Field(default_factory=list)
    high_count: int = 0
    medium_count: int = 0
    low_count: int = 0
    passed: bool = True          # True when no HIGH issues
    threshold_score: int = 100   # 100 minus deductions (see _score)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _slide_title_text(slide) -> str:
    """Return the title placeholder text (empty string if absent)."""
    try:
        return slide.shapes.title.text_frame.text if slide.shapes.title else ""
    except Exception:  # noqa: BLE001
        return ""


def _is_section_slide(slide) -> bool:
    try:
        layout_name = slide.slide_layout.name
        return layout_name in _SECTION_LAYOUTS
    except Exception:  # noqa: BLE001
        return False


def _ph_body_paragraphs(slide) -> tuple[list, bool]:
    """Return (paragraph list, has_body) for the body placeholder of a slide."""
    for shape in slide.shapes:
        try:
            if shape.has_text_frame and shape.placeholder_format is not None:
                ph_idx = shape.placeholder_format.idx
                if ph_idx not in (0, 1):  # 0=title, 1=body (idx 1 is body)
                    continue
                if ph_idx == 1:
                    return list(shape.text_frame.paragraphs), True
        except Exception:  # noqa: BLE001
            pass
    # fallback: second largest text placeholder
    for shape in slide.shapes:
        try:
            if shape.has_text_frame and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 1:
                    return list(shape.text_frame.paragraphs), True
        except Exception:  # noqa: BLE001
            pass
    return [], False


def _score(issues: list[SlideIssue]) -> int:
    """Rule-based threshold score: 100 minus deductions per issue type/severity."""
    deduct = 0
    for iss in issues:
        if iss.severity == "high":
            deduct += 10
        elif iss.severity == "medium":
            deduct += 3
        else:
            deduct += 1
    return max(0, 100 - deduct)


# ---------------------------------------------------------------------------
# Core audit
# ---------------------------------------------------------------------------

def audit_pptx_deterministic(pptx_path: str | Path) -> DeckVisualAuditResult:
    """Inspect every slide in *pptx_path* and return a structured audit result.

    Uses the python-pptx XML model only — no rendering, no network, no LLM.
    """
    from pptx import Presentation  # type: ignore
    from pptx.util import Pt  # type: ignore

    pptx_path = Path(pptx_path)
    issues: list[SlideIssue] = []

    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:  # noqa: BLE001
        return DeckVisualAuditResult(
            pptx_path=str(pptx_path),
            slide_count=0,
            issues=[SlideIssue(
                slide_idx=0, slide_title="",
                issue_type="title_too_long",  # sentinel
                severity="high",
                detail=f"Could not open PPTX: {exc}",
                patch_hint=None,
            )],
            high_count=1, passed=False, threshold_score=0,
        )

    slide_count = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        title_text = _slide_title_text(slide)
        is_section = _is_section_slide(slide)

        # --- title_too_long -----------------------------------------------
        if title_text and len(title_text) > _MAX_TITLE_LEN:
            issues.append(SlideIssue(
                slide_idx=idx, slide_title=title_text,
                issue_type="title_too_long", severity="medium",
                detail=f"Title has {len(title_text)} chars (max {_MAX_TITLE_LEN}).",
                patch_hint=f"Truncate to {_MAX_TITLE_LEN} chars.",
            ))

        # --- too_many_bullets (body placeholder) ---------------------------
        if not is_section:
            body_paras, has_body = _ph_body_paragraphs(slide)
            non_empty_paras = [p for p in body_paras if p.text.strip()]
            if has_body and len(non_empty_paras) > _MAX_BULLETS:
                issues.append(SlideIssue(
                    slide_idx=idx, slide_title=title_text,
                    issue_type="too_many_bullets", severity="medium",
                    detail=f"Body has {len(non_empty_paras)} paragraphs (max {_MAX_BULLETS}).",
                    patch_hint=f"Truncate to {_MAX_BULLETS} paragraphs.",
                ))
            elif not has_body and not is_section:
                # Check if there's any text-bearing shape at all
                has_text = any(
                    shape.has_text_frame and shape.text_frame.text.strip()
                    for shape in slide.shapes
                )
                if not has_text:
                    # Check for image/table — those are OK
                    has_content = any(
                        shape.shape_type in (13, 19)  # MSO_SHAPE_TYPE.PICTURE, TABLE
                        for shape in slide.shapes
                    )
                    if not has_content:
                        issues.append(SlideIssue(
                            slide_idx=idx, slide_title=title_text,
                            issue_type="empty_body", severity="low",
                            detail="Content slide has no body text, image, or table.",
                            patch_hint="Add content or convert to a section slide.",
                        ))

        # --- table_overflow -----------------------------------------------
        for shape in slide.shapes:
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            if len(cell_text) > _MAX_CELL_LEN:
                                issues.append(SlideIssue(
                                    slide_idx=idx, slide_title=title_text,
                                    issue_type="table_overflow", severity="medium",
                                    detail=f"Table cell has {len(cell_text)} chars (max {_MAX_CELL_LEN}).",
                                    patch_hint=f"Clip cell text to {_MAX_CELL_LEN} chars.",
                                ))
            except Exception:  # noqa: BLE001
                pass

        # --- tiny_font + font_drift (run-level scan) ----------------------
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # tiny_font
                        if run.font.size is not None:
                            pt = run.font.size.pt
                            if pt < _MIN_FONT_PT:
                                issues.append(SlideIssue(
                                    slide_idx=idx, slide_title=title_text,
                                    issue_type="tiny_font", severity="high",
                                    detail=f"Run '{run.text[:30]}' has font size {pt:.1f}pt (min {_MIN_FONT_PT}pt).",
                                    patch_hint="Increase font size to at least 10pt.",
                                ))
                        # font_drift (only flag when explicitly set, not inherited)
                        if run.font.name and run.font.name not in _BNK_FONTS:
                            issues.append(SlideIssue(
                                slide_idx=idx, slide_title=title_text,
                                issue_type="font_drift", severity="low",
                                detail=f"Run uses '{run.font.name}' (expected Calibri / Calibri Light).",
                                patch_hint="Reset font to Calibri.",
                            ))
            except Exception:  # noqa: BLE001
                pass

    high = sum(1 for i in issues if i.severity == "high")
    medium = sum(1 for i in issues if i.severity == "medium")
    low = sum(1 for i in issues if i.severity == "low")
    return DeckVisualAuditResult(
        pptx_path=str(pptx_path),
        slide_count=slide_count,
        issues=issues,
        high_count=high,
        medium_count=medium,
        low_count=low,
        passed=high == 0,
        threshold_score=_score(issues),
    )


# ---------------------------------------------------------------------------
# Targeted patcher
# ---------------------------------------------------------------------------

def patch_pptx_overflow(pptx_path: str | Path, issues: list[SlideIssue]) -> str:
    """Apply surgical text fixes for layout issues; write *out_patched.pptx*.

    Handles:
      * title_too_long   — truncate title text with ellipsis
      * too_many_bullets — keep first MAX_BULLETS paragraphs, append "…" marker
      * table_overflow   — clip each overlong cell

    Returns the path to the patched file.
    """
    from pptx import Presentation  # type: ignore

    pptx_path = Path(pptx_path)
    patched_path = pptx_path.parent / "out_patched.pptx"
    shutil.copy2(pptx_path, patched_path)

    prs = Presentation(str(patched_path))
    slides = list(prs.slides)

    # Group issues by slide_idx for efficient processing
    by_slide: dict[int, list[SlideIssue]] = {}
    for iss in issues:
        by_slide.setdefault(iss.slide_idx, []).append(iss)

    for slide_idx, slide_issues in by_slide.items():
        if slide_idx >= len(slides):
            continue
        slide = slides[slide_idx]
        types = {iss.issue_type for iss in slide_issues}

        # Patch title
        if "title_too_long" in types and slide.shapes.title:
            try:
                tf = slide.shapes.title.text_frame
                full_text = tf.text
                if len(full_text) > _MAX_TITLE_LEN:
                    # Modify only the first run of the first paragraph to avoid
                    # disrupting formatting; clear the rest.
                    paras = tf.paragraphs
                    if paras and paras[0].runs:
                        paras[0].runs[0].text = full_text[: _MAX_TITLE_LEN - 1] + "…"
                        for run in paras[0].runs[1:]:
                            run.text = ""
                        for para in paras[1:]:
                            for run in para.runs:
                                run.text = ""
            except Exception:  # noqa: BLE001
                pass

        # Patch body bullets
        if "too_many_bullets" in types:
            try:
                _, has_body = _ph_body_paragraphs(slide)
                if has_body:
                    for shape in slide.shapes:
                        if (shape.has_text_frame
                                and shape.placeholder_format is not None
                                and shape.placeholder_format.idx == 1):
                            paras = shape.text_frame.paragraphs
                            non_empty = [p for p in paras if p.text.strip()]
                            if len(non_empty) > _MAX_BULLETS:
                                # Blank out paragraphs beyond MAX_BULLETS
                                for para in non_empty[_MAX_BULLETS:]:
                                    for run in para.runs:
                                        run.text = ""
                                # Append ellipsis to last kept paragraph
                                last_kept = non_empty[_MAX_BULLETS - 1]
                                if last_kept.runs:
                                    last_kept.runs[-1].text += " …"
                            break
            except Exception:  # noqa: BLE001
                pass

        # Patch table cells
        if "table_overflow" in types:
            for shape in slide.shapes:
                try:
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if not cell.text_frame:
                                    continue
                                cell_text = cell.text_frame.text
                                if len(cell_text) > _MAX_CELL_LEN:
                                    paras = cell.text_frame.paragraphs
                                    if paras and paras[0].runs:
                                        paras[0].runs[0].text = cell_text[: _MAX_CELL_LEN - 1] + "…"
                                        for run in paras[0].runs[1:]:
                                            run.text = ""
                                        for para in paras[1:]:
                                            for run in para.runs:
                                                run.text = ""
                except Exception:  # noqa: BLE001
                    pass

    prs.save(str(patched_path))
    return str(patched_path)


# ---------------------------------------------------------------------------
# Persistence helpers
# ---------------------------------------------------------------------------

def write_visual_audit(result: DeckVisualAuditResult, workspace: Path) -> None:
    path = Path(workspace) / VISUAL_AUDIT_NAME
    path.write_text(
        json.dumps(result.model_dump(), indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def load_visual_audit(workspace: Path) -> DeckVisualAuditResult | None:
    path = Path(workspace) / VISUAL_AUDIT_NAME
    if not path.exists():
        return None
    try:
        return DeckVisualAuditResult.model_validate(
            json.loads(path.read_text(encoding="utf-8"))
        )
    except Exception:  # noqa: BLE001
        return None


# ---------------------------------------------------------------------------
# Human-readable summary
# ---------------------------------------------------------------------------

def format_visual_audit(result: DeckVisualAuditResult) -> str:
    status = "PASSED" if result.passed else "ISSUES FOUND"
    lines = [
        f"\nDECK VISUAL AUDIT [{status}] — {result.slide_count} slides, "
        f"score {result.threshold_score}/100 "
        f"(HIGH:{result.high_count} MED:{result.medium_count} LOW:{result.low_count})"
    ]
    if not result.issues:
        lines.append("  No layout issues detected.")
    else:
        for iss in result.issues:
            icon = {"high": "🔴", "medium": "🟡", "low": "⚪"}.get(iss.severity, "")
            lines.append(
                f"  {icon} slide {iss.slide_idx + 1} [{iss.issue_type}]: {iss.detail}"
            )
            if iss.patch_hint:
                lines.append(f"      → {iss.patch_hint}")
    return "\n".join(lines)
