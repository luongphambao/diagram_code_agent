"""Editable BnK PowerPoint proposal generation for architecture diagrams."""

from __future__ import annotations

import datetime as dt
import json
import re
import warnings
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt

# --- BnK brand palette (Calibri + corporate blue, from template style guide) ---
BNK_BLUE = RGBColor(0x1F, 0x4E, 0x78)   # primary corporate blue (header rows, dividers)
BNK_CYAN = RGBColor(0x00, 0x9F, 0xDF)   # secondary accent
BNK_LIGHT = RGBColor(0xE9, 0xF0, 0xF7)  # light tint for alternating table rows
BNK_ACCENT = RGBColor(0xC0, 0x3A, 0x2B)  # red accent (required/important)
BNK_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BNK_TEXT = RGBColor(0x33, 0x33, 0x33)
BNK_FONT = "Calibri"

from reporting import (
    DEFAULT_REPORT_SECTIONS,
    assemble_report_data,
    normalize_sections,
    read_json_file,
    record_artifact_inventory,
    record_report_step,
)

DEFAULT_PPT_SECTIONS = [
    "cover",
    "executive_summary",
    "solution_overview",
    "scope",
    "architecture_diagram",
    "technical_stack",
    "key_decisions",
    "delivery_plan",
    "pricing",
    "risks",
    "appendix",
]

SECTION_ALIASES = {
    "solution": "solution_overview",
    "overview": "solution_overview",
    "diagram": "architecture_diagram",
    "architecture": "architecture_diagram",
    "techstack": "technical_stack",
    "tech_stack": "technical_stack",
    "decisions": "key_decisions",
    "delivery": "delivery_plan",
    "team": "delivery_plan",
    "wbs": "delivery_plan",
    "price": "pricing",
    "pricing": "pricing",
    "capex": "pricing",
    "cost": "pricing",
    "risk": "risks",
    "artifact": "appendix",
    "reference": "appendix",
}

# Layout names below MUST match the actual layout names inside the BnK template
# (verified via python-pptx). Note the double space in the separator layout name.
VALID_LAYOUTS = frozenset({
    "Cover-01",
    "Head Page",
    "Head-01",
    "Detail-01",
    "Overview-01",
    "Empty",
    "BnK",
    "C2 -  Separator/ Dark",
})

# Closing/separator layouts that are appended via _append_thank_you, never rendered inline.
CLOSING_LAYOUTS = frozenset({"BnK", "C2 -  Separator/ Dark"})

# Structured content blocks an outline slide may carry (in addition to plain bullets).
VALID_BLOCKS = frozenset({
    "bullets",
    "tech_stack_table",
    "func_nfr",
    "sdlc",
    "delivery_effort",
    "pricing",
    "milestones",
    "team",
    "gantt",  # rendered via _gantt_slide; not yet dispatched from _render_block's
              # legacy outline path — see deck_resolver._b_master_plan for the params shape.
})

OUTLINE_TARGET_MIN = 20
OUTLINE_TARGET_MAX = 30

_OUTLINE_SYSTEM_PROMPT = """\
You are a senior solution architect at BnK, a Vietnamese technology consultancy.
Generate a professional PowerPoint proposal slide outline for a client project.

OUTPUT FORMAT — each slide MUST have these fields:
{
  "title": "SECTION | Sub-topic",
  "layout": "<layout_name>",
  "block": "<block_type>",
  "bullets": ["...", "..."],
  "asset_ref": null
}

VALID LAYOUT NAMES (use EXACTLY as written):
  "Cover-01"              — Opening cover. Use ONCE as slide 1. Set title to "".
  "Head Page"            — Major section divider with Roman numeral (I., II., III., …).
  "Head-01"              — Secondary section header, no Roman numeral.
  "Detail-01"            — Content slide with title + bullets. Most common.
  "Overview-01"          — Overview: put one subtitle string in bullets[0], no other bullets.
  "Empty"                — Full-width image slide. Use with asset_ref: "architecture_diagram".
  "BnK"                  — Closing brand slide. Appended automatically; do NOT emit.
  "C2 -  Separator/ Dark" — Dark separator. Appended automatically; do NOT emit.

CONTENT BLOCKS — set "block" to render a professional table instead of bullets.
For these, "layout" is ignored (always Detail-01) and "bullets" may be []:
  "bullets"          — (default) plain bullet slide on the given layout.
  "tech_stack_table" — table Layer | Technology | Description, built from TECH_STACK data.
  "func_nfr"         — two columns: Functional vs Non-Functional Requirements.
  "sdlc"             — SCOPE OF WORK SDLC phase table (Analysis→Maintenance).
  "delivery_effort"  — effort table Code | Module | MD, built from WBS_SUMMARY.
  "pricing"          — CAPEX cost table (NET, excluding taxes).
  "milestones"       — payment milestones table (30/30/30/10).
  "team"             — Client Team vs BnK Team table.
When using a block, DO NOT also invent bullet content for that table — leave bullets [].

SLIDE COUNT: Generate exactly 20-30 slides total.

TYPICAL STRUCTURE (adapt to the actual project data; use blocks where noted):
  1.   Cover-01                          cover
  2.   Head Page                         "I. Executive Summary"
  3-4. Detail-01 (bullets)               executive highlights, business value
  5.   Head Page                         "II. Proposed Solution"
  6.   Overview-01                       solution scope overview (subtitle in bullets[0])
  7.   block:"func_nfr"                  "PROPOSED SOLUTION | Requirements"
  8-9. Detail-01 (bullets)               approach, key features
  10.  Head-01                           "Architecture Overview"
  11.  Empty (asset_ref diagram)         architecture diagram
  12.  block:"tech_stack_table"          "PROPOSED SOLUTION | Technical Stack"
  13.  Head Page                         "IV. Scope of Work"
  14.  block:"sdlc"                      "SCOPE OF WORK | SDLC Phases"
  15.  Detail-01 (bullets)               deliverables, assumptions, change request
  16.  Head Page                         "V. Project Delivery"
  17.  block:"delivery_effort"           "PROJECT DELIVERY | Estimated Effort"
  18.  block:"team"                      "PROJECT DELIVERY | Team Structure"
  19.  Detail-01 (bullets)               risks & mitigations
  20.  Head Page                         "VI. Pricing"
  21.  block:"pricing"                   "PRICING | CAPEX"
  22.  block:"milestones"                "PRICING | Payment Milestones"
  (BnK closing slide is appended automatically — do not include it.)

TITLE FORMAT:
  - Detail-01 / blocks: "SECTION | Sub-topic"  e.g. "PROPOSED SOLUTION | Technical Stack"
  - Head Page / Head-01: section name with Roman numeral  e.g. "IV. Scope of Work"
  - Cover-01: empty string ""

RULES:
  1. Use ONLY the layout names and block names listed — no other values.
  2. Cover-01 must be slide 1; do NOT emit BnK / separator (appended automatically).
  3. Include the Empty/diagram slide ONLY if HAS_ARCHITECTURE_DIAGRAM is yes.
  4. Include block:"delivery_effort" only if WBS_SUMMARY has data; otherwise use bullets.
  5. Base ALL content on the actual project data — no generic placeholders.
  6. OUTPUT: Return ONLY the JSON array. No code fences, no text outside the array.\
"""


class PPTProposalError(RuntimeError):
    """Raised when PPT proposal generation cannot complete."""


def _template_path() -> Path:
    # parents[0]=src/, parents[1]=backend/ — bundled alongside the package
    return Path(__file__).resolve().parents[1] / "templates" / "bnk_proposal_template.pptx"


def _clip(text: Any, limit: int = 220) -> str:
    value = " ".join(str(text or "").split())
    return value[:limit].rstrip() + ("..." if len(value) > limit else "")


def _as_list(value: Any) -> list[Any]:
    if isinstance(value, list):
        return value
    if value in (None, ""):
        return []
    return [value]


def normalize_ppt_sections(sections: list[str] | None) -> tuple[list[str], list[str]]:
    if not sections:
        return DEFAULT_PPT_SECTIONS.copy(), []
    out: list[str] = []
    unrecognized: list[str] = []
    for raw in sections:
        name = SECTION_ALIASES.get(str(raw).strip().lower(), str(raw).strip().lower())
        if name in DEFAULT_PPT_SECTIONS and name not in out:
            out.append(name)
        elif name not in DEFAULT_PPT_SECTIONS:
            unrecognized.append(str(raw).strip())
    if unrecognized and len(unrecognized) > len(sections) / 2:
        return DEFAULT_PPT_SECTIONS.copy(), unrecognized
    return out or DEFAULT_PPT_SECTIONS.copy(), unrecognized


_ROMAN = [(1000,"M"),(900,"CM"),(500,"D"),(400,"CD"),(100,"C"),(90,"XC"),
          (50,"L"),(40,"XL"),(10,"X"),(9,"IX"),(5,"V"),(4,"IV"),(1,"I")]


def _roman(n: int) -> str:
    result = ""
    for val, numeral in _ROMAN:
        while n >= val:
            result += numeral
            n -= val
    return result


def _layout(prs: Presentation, *names: str):
    wanted = [n.lower() for n in names]
    by_name = {layout.name.lower(): layout for layout in prs.slide_layouts}
    for name in wanted:  # honour caller priority order
        if name in by_name:
            return by_name[name]
    # Prefer a neutral content/blank layout over slide_layouts[0] (which is the cover).
    for safe in ("detail-01", "blank", "empty"):
        if safe in by_name:
            warnings.warn(
                f"PPT layout(s) {list(names)!r} not found; using {by_name[safe].name!r} fallback.",
                stacklevel=3,
            )
            return by_name[safe]
    warnings.warn(
        f"PPT layout(s) {list(names)!r} not found in template; using layout[0] as fallback.",
        stacklevel=3,
    )
    return prs.slide_layouts[0]


def _clear_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public clear API.
    for slide_id in list(slide_id_list):
        rid = slide_id.rId
        prs.part.drop_rel(rid)
        slide_id_list.remove(slide_id)


def _set_placeholder_text(slide, idx: int, text: str) -> bool:
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx == idx:
                shape.text = text
                return True
        except Exception:
            continue
    return False


def _body_placeholder(slide):
    """Return the slide's BODY/content placeholder (idx 13 in Detail-01) if present."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    for shape in slide.placeholders:
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            continue
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def _fill_bullets_placeholder(slide, items: list[Any], *, font_size: int = 16, limit: int = 7) -> bool:
    """Fill the template's body placeholder with bullets (inherits brand styling).

    Returns True if a placeholder was used, False if none exists.
    """
    ph = _body_placeholder(slide)
    if ph is None:
        return False
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    bullets = [_clip(item, 180) for item in items[:limit] if str(item or "").strip()]
    if not bullets:
        bullets = ["Details will be confirmed during proposal review."]
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return True


def _add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 16,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold


def _add_bullets(
    slide,
    items: list[Any],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 16,
    limit: int = 7,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullets = [_clip(item, 180) for item in items[:limit] if str(item or "").strip()]
    if not bullets:
        bullets = ["Details will be confirmed during proposal review."]
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)


def _add_title(slide, title: str) -> None:
    if not _set_placeholder_text(slide, 0, title):
        _add_textbox(slide, title, 0.65, 0.35, 12.0, 0.55, font_size=28, bold=True)


def _add_footer(slide, slide_no: int) -> None:
    _add_textbox(slide, str(slide_no), 12.25, 6.9, 0.45, 0.2, font_size=8, align=PP_ALIGN.RIGHT)


def _image_fit(slide, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    from PIL import Image

    with Image.open(image_path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih if ih else box_ratio
    if img_ratio > box_ratio:
        final_w = w
        final_h = w / img_ratio
    else:
        final_h = h
        final_w = h * img_ratio
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))


def _clone_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    blank = _layout(prs, "Blank", "Empty")
    slide = prs.slides.add_slide(blank)
    for shape in source.shapes:
        slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")  # noqa: SLF001
    return slide


def _section_slide(prs: Presentation, title: str, slide_no: int):
    slide = prs.slides.add_slide(_layout(prs, "Head Page", "Head-01"))
    _add_title(slide, title)
    _add_footer(slide, slide_no)
    return slide


def _detail_slide(prs: Presentation, title: str, bullets: list[Any], slide_no: int):
    slide = prs.slides.add_slide(_layout(prs, "Detail-01"))
    _add_title(slide, title)
    # Prefer the template's designed body placeholder; fall back to a positioned textbox.
    if not _fill_bullets_placeholder(slide, bullets):
        _add_bullets(slide, bullets, 0.85, 1.35, 11.6, 4.95)
    _add_footer(slide, slide_no)
    return slide


def _overview_slide(prs: Presentation, title: str, subtitle: str, slide_no: int):
    slide = prs.slides.add_slide(_layout(prs, "Overview-01"))
    _set_placeholder_text(slide, 0, title)
    _set_placeholder_text(slide, 1, subtitle)
    _add_footer(slide, slide_no)
    return slide


def _cover_slide(prs: Presentation, report: dict[str, Any], slide_no: int):
    slide = prs.slides.add_slide(_layout(prs, "Cover-01"))
    _set_placeholder_text(slide, 0, report["title"])
    _set_placeholder_text(slide, 1, report["subtitle"])
    date_text = dt.datetime.now().strftime("%B %d, %Y")
    _add_textbox(slide, date_text, 0.75, 6.45, 4.4, 0.3, font_size=13)
    if report.get("brand"):
        _add_textbox(slide, str(report["brand"]), 0.75, 5.95, 4.4, 0.3, font_size=12)
    _add_footer(slide, slide_no)
    return slide


def _diagram_slide(prs: Presentation, report: dict[str, Any], workspace: Path, slide_no: int):
    slide = prs.slides.add_slide(_layout(prs, "Blank", "Empty"))
    _add_title(slide, report["blueprint"].get("diagram_title") or "Application Architecture")
    diagram = workspace / "out.body.png"
    if not diagram.exists():
        diagram = workspace / "out.png"
    if diagram.exists():
        _image_fit(slide, diagram, 0.55, 1.05, 12.15, 5.55)
    else:
        _add_textbox(slide, "No diagram image is available.", 1.0, 2.8, 11.0, 0.4, font_size=16)
    _add_footer(slide, slide_no)
    return slide


# --------------------------------------------------------------------------- #
# Native (editable) brand-styled tables and the BnK-specific slide builders.
# --------------------------------------------------------------------------- #

_CONTENT_X = 0.6
_CONTENT_Y = 1.35
_CONTENT_W = 12.1
_CONTENT_H = 5.35


def _style_cell(cell, text: str, *, fill: RGBColor, color: RGBColor, bold: bool, size: int, align=PP_ALIGN.LEFT) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = _clip(text, 240)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = BNK_FONT
        run.font.color.rgb = color


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[Any]],
    *,
    x: float = _CONTENT_X,
    y: float = _CONTENT_Y,
    w: float = _CONTENT_W,
    h: float = _CONTENT_H,
    col_widths: list[float] | None = None,
    header_size: int = 13,
    body_size: int = 11,
    max_rows: int = 12,
):
    """Add a brand-styled, editable table. Header row = BnK blue; alternate body tint."""
    rows = [r for r in rows if any(str(c or "").strip() for c in r)][:max_rows]
    if not rows:
        rows = [["—"] * len(headers)]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = gfx.table
    table.first_row = False  # disable theme banding; we colour manually
    table.horz_banding = False
    if col_widths and len(col_widths) == n_cols:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, head in enumerate(headers):
        _style_cell(table.cell(0, c), str(head), fill=BNK_BLUE, color=BNK_WHITE, bold=True, size=header_size)
    for r, row in enumerate(rows, start=1):
        fill = BNK_LIGHT if r % 2 == 0 else BNK_WHITE
        for c in range(n_cols):
            value = row[c] if c < len(row) else ""
            _style_cell(table.cell(r, c), str(value), fill=fill, color=BNK_TEXT, bold=False, size=body_size)
    return table


def _table_slide(prs: Presentation, title: str, headers, rows, slide_no: int, **kw):
    """A Detail-01 slide whose content area holds a brand-styled table."""
    slide = prs.slides.add_slide(_layout(prs, "Detail-01"))
    _add_title(slide, title)
    _add_table(slide, headers, rows, **kw)
    _add_footer(slide, slide_no)
    return slide


# Top-level tech_stack scalar keys that _tech_items may surface as bogus "layers".
_TECH_META_KEYS = frozenset({
    "estimated_total_monthly_cost_usd", "assumptions", "scaling_roadmap", "notes", "summary",
})


def _tech_stack_table_slide(prs: Presentation, report: dict[str, Any], slide_no: int, title: str = "PROPOSED SOLUTION | Technical Stack"):
    rows = []
    for item in report.get("tech_items", [])[:12]:
        layer = item.get("layer") or "Layer"
        if str(layer).strip().lower() in _TECH_META_KEYS:
            continue
        choice = item.get("choice") or item.get("name") or "TBD"
        rationale = _clip(item.get("rationale") or "", 140)
        rows.append([layer, choice, rationale])
    if not rows:
        rows = [["Frontend", "TBD", ""], ["Backend", "TBD", ""], ["Database", "TBD", ""]]
    return _table_slide(
        prs, title, ["Layer", "Technology", "Description"], rows, slide_no,
        col_widths=[2.4, 3.2, 6.5],
    )


def _functional_nfr_slide(prs: Presentation, report: dict[str, Any], slide_no: int, title: str = "PROPOSED SOLUTION | Requirements"):
    brief = report.get("brief") or {}
    func = _as_list(brief.get("functional_requirements"))[:8]
    nfr = _as_list(brief.get("non_functional_requirements"))[:8]
    slide = prs.slides.add_slide(_layout(prs, "Detail-01"))
    _add_title(slide, title)
    half_w = (_CONTENT_W - 0.4) / 2
    _add_table(
        slide, ["Functional Requirements"], [[r] for r in func] or [["To be confirmed"]],
        x=_CONTENT_X, y=_CONTENT_Y, w=half_w, h=_CONTENT_H, header_size=13, body_size=11,
    )
    _add_table(
        slide, ["Non-Functional Requirements"], [[r] for r in nfr] or [["To be confirmed"]],
        x=_CONTENT_X + half_w + 0.4, y=_CONTENT_Y, w=half_w, h=_CONTENT_H, header_size=13, body_size=11,
    )
    _add_footer(slide, slide_no)
    return slide


_SDLC_DEFAULT = [
    ("Analysis", "Collect & validate requirements with client; wireframing", "BRD, Wireframes"),
    ("Design", "System architecture and UI/UX design", "System Design, UI/UX"),
    ("Development", "Implementation in agile sprints", "Source Code"),
    ("Testing", "Manual QC, SIT & UAT support, load & security test", "Test Report"),
    ("Deployment", "Deploy to DEV / UAT / PROD environments", "Deployment Guide"),
    ("Maintenance", "Post go-live support to fix defects", "Support"),
]


def _sdlc_scope_slide(prs: Presentation, report: dict[str, Any], workspace: Path, slide_no: int, title: str = "SCOPE OF WORK | SDLC Phases"):
    rows = [[p, d, dl] for (p, d, dl) in _SDLC_DEFAULT]
    return _table_slide(
        prs, title, ["Phase", "Activities", "Deliverables"], rows, slide_no,
        col_widths=[2.2, 6.4, 3.5],
    )


def _delivery_effort_slide(prs: Presentation, workspace: Path, slide_no: int, title: str = "PROJECT DELIVERY | Estimated Effort"):
    wbs = read_json_file(workspace / "wbs.json", {})
    headers = ["Code", "Module", "Effort (MD)"]
    rows: list[list[Any]] = []
    if isinstance(wbs, dict) and wbs:
        for mod in _as_list(wbs.get("effort_by_module"))[:10]:
            if isinstance(mod, dict):
                rows.append([mod.get("code", ""), mod.get("name", "Module"), mod.get("total_md", 0)])
        totals = wbs.get("effort_totals") or {}
        if totals:
            rows.append(["", "TOTAL", totals.get("total_mandays", 0)])
    if not rows:
        rows = [["", "Effort will be finalized after WBS approval.", ""]]
    return _table_slide(
        prs, title, headers, rows, slide_no, col_widths=[1.6, 7.5, 3.0],
    )


def _pricing_slide(prs: Presentation, report: dict[str, Any], slide_no: int, title: str = "PRICING | CAPEX"):
    total = report.get("tech_total_cost")
    rows: list[list[Any]] = []
    for item in report.get("tech_items", [])[:8]:
        layer = item.get("layer") or "Item"
        if str(layer).strip().lower() in _TECH_META_KEYS:
            continue
        choice = item.get("choice") or item.get("name") or ""
        rows.append([f"{layer}: {choice}".strip(": "), "—"])
    rows.append(["Total (NET, excluding taxes/VAT)", f"{total} USD" if total else "XXX USD"])
    return _table_slide(
        prs, title, ["Cost Item", "Amount"], rows, slide_no, col_widths=[8.5, 3.6],
    )


_DEFAULT_MILESTONES = [
    ("1", "Contract sign-off / Analysis complete", "30%"),
    ("2", "Completion of Development", "30%"),
    ("3", "Completion of UAT", "30%"),
    ("4", "Completion of Nursing Period", "10%"),
]


def _payment_milestones_slide(prs: Presentation, slide_no: int, title: str = "PRICING | Payment Milestones"):
    rows = [[n, name, pct] for (n, name, pct) in _DEFAULT_MILESTONES]
    return _table_slide(
        prs, title, ["#", "Milestone", "Payment"], rows, slide_no, col_widths=[1.0, 8.6, 2.5],
    )


def _gantt_slide(
    prs: Presentation, params: dict[str, Any], slide_no: int,
    title: str = "PROJECT DELIVERY | Master Plan & Milestones",
):
    """The Master Plan Gantt — same schedule as the WBS Excel '3. Delivery Plan' sheet
    (deck_resolver._b_master_plan calls wbs_excel._module_schedule so the two never drift).

    Column granularity adapts to project length so the grid stays legible on a 13.3" slide:
    weekly for short projects, sprint-grouped, then monthly for long ones — mirroring how
    the Excel sheet groups Month > Sprint > Week headers, just collapsed to one row.
    """
    weeks = int(params.get("weeks") or 0)
    months = int(params.get("months") or 0)
    sprints = int(params.get("sprints") or 0)
    rows = params.get("gantt_rows") or []

    if weeks <= 20 and weeks:
        n_cols, weeks_per_col, label = weeks, 1, lambda i: f"W{i + 1}"
    elif weeks <= 40 and sprints:
        n_cols, weeks_per_col, label = sprints, 2, lambda i: f"S{i + 1}"
    else:
        n_cols, weeks_per_col, label = (months or 1), 4, lambda i: f"M{i + 1}"

    slide = prs.slides.add_slide(_layout(prs, "Detail-01"))
    _add_title(slide, title)
    if not rows or not n_cols:
        _add_textbox(slide, "Delivery timeline not yet available.", 0.85, 1.6, 11.0, 0.4, font_size=14)
        _add_footer(slide, slide_no)
        return slide

    name_w = 3.0
    grid_w = _CONTENT_W - name_w
    col_w = grid_w / n_cols
    n_rows = min(len(rows), 16) + 1
    gfx = slide.shapes.add_table(
        n_rows, n_cols + 1, Inches(_CONTENT_X), Inches(_CONTENT_Y), Inches(_CONTENT_W), Inches(_CONTENT_H)
    )
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(name_w)
    for i in range(n_cols):
        table.columns[i + 1].width = Inches(col_w)

    _style_cell(table.cell(0, 0), "Module", fill=BNK_BLUE, color=BNK_WHITE, bold=True, size=10)
    for i in range(n_cols):
        _style_cell(table.cell(0, i + 1), label(i), fill=BNK_BLUE, color=BNK_WHITE, bold=True,
                    size=7, align=PP_ALIGN.CENTER)

    for r, m in enumerate(rows[:16], start=1):
        _style_cell(table.cell(r, 0), f"{m.get('code', '')} {m.get('name', '')}".strip(),
                   fill=BNK_WHITE, color=BNK_TEXT, bold=False, size=9)
        start_col = max(0, (int(m.get("start_week", 1)) - 1) // weeks_per_col)
        end_col = min(n_cols - 1, (int(m.get("end_week", 1)) - 1) // weeks_per_col)
        for i in range(n_cols):
            active = start_col <= i <= end_col
            _style_cell(table.cell(r, i + 1), "", fill=BNK_CYAN if active else BNK_WHITE,
                       color=BNK_WHITE, bold=False, size=1)

    _add_footer(slide, slide_no)
    return slide


def _team_slide(prs: Presentation, report: dict[str, Any], slide_no: int, title: str = "PROJECT DELIVERY | Team Structure"):
    rows = [
        ["Technical Lead", "Technical Lead"],
        ["Business Analyst", "Developer(s)"],
        ["Project Manager", "BA / Tester"],
        ["", "Project Manager"],
    ]
    return _table_slide(
        prs, title, ["Client Team", "BnK Delivery Team"], rows, slide_no, col_widths=[6.0, 6.1],
    )


def _tech_bullets(report: dict[str, Any]) -> list[str]:
    items = []
    for item in report.get("tech_items", [])[:8]:
        layer = item.get("layer") or "Layer"
        choice = item.get("choice") or item.get("name") or "TBD"
        rationale = item.get("rationale") or ""
        items.append(f"{layer}: {choice}" + (f" - {_clip(rationale, 100)}" if rationale else ""))
    return items


def _delivery_bullets(workspace: Path) -> list[str]:
    wbs = read_json_file(workspace / "wbs.json", {})
    if not isinstance(wbs, dict) or not wbs:
        return [
            "Delivery plan will be finalized after WBS approval.",
            "Recommended phases: discovery, solution design, implementation, testing, UAT, launch support.",
        ]
    totals = wbs.get("effort_totals") or {}
    timeline = wbs.get("timeline") or {}
    bullets = [
        f"Total effort: {totals.get('total_mandays', 0)} MD / {totals.get('total_manmonths', 0)} MM.",
        f"Timeline: {timeline.get('weeks', 0)} weeks / {timeline.get('months', 0)} months.",
    ]
    for module in _as_list(wbs.get("effort_by_module"))[:5]:
        if isinstance(module, dict):
            bullets.append(f"{module.get('code', '')} {module.get('name', 'Module')}: {module.get('total_md', 0)} MD.")
    return bullets


def _build_outline_context(report: dict[str, Any], workspace: Path) -> str:
    """Serialize report data to compact plain-text for the LLM outline prompt."""
    lines: list[str] = []

    lines.append(f"TITLE: {_clip(report.get('title', ''), 120)}")
    lines.append(f"SUBTITLE: {_clip(report.get('subtitle', ''), 120)}")
    lines.append(f"BRAND: {_clip(report.get('brand', ''), 80)}")

    brief = report.get("brief") or {}
    lines.append(f"\nOBJECTIVE: {_clip(brief.get('objective', ''), 300)}")
    stakeholders = _as_list(brief.get("stakeholders"))[:5]
    if stakeholders:
        lines.append(f"STAKEHOLDERS: {', '.join(_clip(str(s), 60) for s in stakeholders)}")

    func_reqs = _as_list(brief.get("functional_requirements"))[:8]
    if func_reqs:
        lines.append("\nFUNCTIONAL_REQUIREMENTS:")
        for r in func_reqs:
            lines.append(f"  - {_clip(r, 200)}")

    nfr = _as_list(brief.get("non_functional_requirements"))[:5]
    if nfr:
        lines.append("\nNON_FUNCTIONAL_REQUIREMENTS:")
        for r in nfr:
            lines.append(f"  - {_clip(r, 200)}")

    blueprint = report.get("blueprint") or {}
    if blueprint.get("pattern"):
        lines.append(f"\nARCHITECTURE_PATTERN: {_clip(blueprint['pattern'], 120)}")
    if report.get("pattern_rationale"):
        lines.append(f"PATTERN_RATIONALE: {_clip(report['pattern_rationale'], 200)}")

    decisions = _as_list(blueprint.get("key_decisions"))[:6]
    if decisions:
        lines.append("\nKEY_DECISIONS:")
        for d in decisions:
            lines.append(f"  - {_clip(d, 200)}")

    tech_items = report.get("tech_items") or []
    if tech_items:
        lines.append("\nTECH_STACK:")
        for item in tech_items[:10]:
            layer = item.get("layer") or "Layer"
            choice = item.get("choice") or item.get("name") or "TBD"
            rationale = _clip(item.get("rationale") or "", 80)
            lines.append(f"  - {layer}: {choice}" + (f" ({rationale})" if rationale else ""))

    exec_pts = _as_list(report.get("executive_points"))[:5]
    if exec_pts:
        lines.append("\nEXECUTIVE_POINTS:")
        for p in exec_pts:
            lines.append(f"  - {_clip(p, 200)}")

    if report.get("business_value"):
        lines.append(f"\nBUSINESS_VALUE: {_clip(report['business_value'], 250)}")
    if report.get("technical_value"):
        lines.append(f"TECHNICAL_VALUE: {_clip(report['technical_value'], 250)}")

    risks = [r for r in _as_list(report.get("risks")) if isinstance(r, dict)][:5]
    if risks:
        lines.append("\nRISKS:")
        for r in risks:
            lines.append(
                f"  - {r.get('type', 'Risk')}: {_clip(r.get('detail', ''), 150)}"
                f" → {_clip(r.get('recommendation', ''), 100)}"
            )

    wbs = read_json_file(workspace / "wbs.json", {})
    if isinstance(wbs, dict) and wbs:
        totals = wbs.get("effort_totals") or {}
        timeline = wbs.get("timeline") or {}
        lines.append("\nWBS_SUMMARY:")
        lines.append(
            f"  Total: {totals.get('total_mandays', 0)} MD"
            f" / {totals.get('total_manmonths', 0)} MM"
        )
        lines.append(
            f"  Timeline: {timeline.get('weeks', 0)} weeks"
            f" / {timeline.get('months', 0)} months"
        )
        for mod in _as_list(wbs.get("effort_by_module"))[:6]:
            if isinstance(mod, dict):
                lines.append(
                    f"  - {mod.get('code', '')} {mod.get('name', 'Module')}:"
                    f" {mod.get('total_md', 0)} MD"
                )
    else:
        lines.append("\nWBS_SUMMARY: Not yet available.")

    has_diagram = (workspace / "out.body.png").exists() or (workspace / "out.png").exists()
    lines.append(f"\nHAS_ARCHITECTURE_DIAGRAM: {'yes' if has_diagram else 'no'}")

    return "\n".join(lines)


def _parse_outline(raw_text: str) -> list[dict[str, Any]] | None:
    """Parse LLM output into a validated pagecontent list. Returns None on any failure."""
    text = raw_text.strip()

    def _try(s: str) -> list | None:
        try:
            obj = json.loads(s)
            return obj if isinstance(obj, list) else None
        except (json.JSONDecodeError, ValueError):
            return None

    result = (
        _try(text)
        or (lambda m: _try(m.group(1)) if m else None)(
            re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", text)
        )
        or (
            (lambda s, e: _try(text[s : e + 1]) if s != -1 and e > s else None)(
                text.find("["), text.rfind("]")
            )
        )
    )

    if not isinstance(result, list):
        return None

    valid: list[dict[str, Any]] = []
    for item in result:
        if not isinstance(item, dict):
            continue
        layout = str(item.get("layout") or "Detail-01")
        if layout not in VALID_LAYOUTS:
            layout = "Detail-01"
        bullets = [
            _clip(b, 200)
            for b in _as_list(item.get("bullets") or [])[:7]
            if str(b or "").strip()
        ]
        block = str(item.get("block") or "bullets")
        if block not in VALID_BLOCKS:
            block = "bullets"
        valid.append(
            {
                "title": _clip(str(item.get("title") or ""), 80),
                "layout": layout,
                "block": block,
                "bullets": bullets,
                "asset_ref": item.get("asset_ref") or None,
            }
        )

    return valid if len(valid) >= 5 else None


def _generate_slide_outline(
    report: dict[str, Any],
    workspace: Path,
    sections: list[str],
) -> list[dict[str, Any]] | None:
    """Call LLM to generate a pagecontent slide outline. Returns None on any failure."""
    try:
        from config import get_model, make_llm
        from langchain_core.messages import HumanMessage, SystemMessage
    except ImportError:
        return None

    try:
        model = get_model("ppt_outline", fallback=get_model("main", "mimo-v2.5"))
        llm = make_llm(model)
    except Exception:
        return None

    context = _build_outline_context(report, workspace)
    sections_note = (
        ""
        if set(sections) == set(DEFAULT_PPT_SECTIONS)
        else (
            f"\nNOTE: Only include slides relevant to these sections: "
            + ", ".join(sections)
            + "."
        )
    )
    user_msg = (
        f"Here is the project data for the proposal:\n\n{context}{sections_note}\n\n"
        "Generate the slide outline now. Return ONLY the JSON array."
    )

    try:
        response = llm.invoke(
            [SystemMessage(content=_OUTLINE_SYSTEM_PROMPT), HumanMessage(content=user_msg)]
        )
        raw = response.content if hasattr(response, "content") else str(response)
        return _parse_outline(raw)
    except Exception:
        return None


def _render_slide(
    prs: Presentation,
    spec: dict[str, Any],
    report: dict[str, Any],
    workspace: Path,
    slide_no: int,
) -> None:
    """Dispatch one slide spec to the appropriate slide builder."""
    layout = spec.get("layout", "Detail-01")
    title = spec.get("title") or ""
    bullets = spec.get("bullets") or []
    asset_ref = spec.get("asset_ref")
    block = spec.get("block") or "bullets"

    # Structured content blocks take precedence over the plain layout dispatch.
    if block != "bullets":
        _render_block(prs, block, title, report, workspace, slide_no)
        return

    if layout == "Cover-01":
        _cover_slide(prs, report, slide_no)
    elif layout == "Head Page":
        slide = prs.slides.add_slide(_layout(prs, "Head Page"))
        _add_title(slide, title)
        _add_footer(slide, slide_no)
    elif layout == "Head-01":
        slide = prs.slides.add_slide(_layout(prs, "Head-01"))
        _add_title(slide, title)
        _add_footer(slide, slide_no)
    elif layout == "Detail-01":
        _detail_slide(prs, title, bullets, slide_no)
    elif layout == "Overview-01":
        _overview_slide(prs, title, bullets[0] if bullets else "", slide_no)
    elif layout == "Empty":
        if asset_ref == "architecture_diagram":
            _diagram_slide(prs, report, workspace, slide_no)
        else:
            slide = prs.slides.add_slide(_layout(prs, "Empty", "Blank"))
            if title:
                _add_title(slide, title)
            _add_footer(slide, slide_no)
    elif layout in CLOSING_LAYOUTS:
        pass  # Closing slides are always appended via _append_thank_you after the loop.
    else:
        _detail_slide(prs, title, bullets, slide_no)


def _render_block(
    prs: Presentation,
    block: str,
    title: str,
    report: dict[str, Any],
    workspace: Path,
    slide_no: int,
) -> None:
    """Render a structured (table-based) BnK slide from a block type."""
    if block == "tech_stack_table":
        _tech_stack_table_slide(prs, report, slide_no, title or "PROPOSED SOLUTION | Technical Stack",
                                workspace=workspace)
    elif block == "func_nfr":
        _functional_nfr_slide(prs, report, slide_no, title or "PROPOSED SOLUTION | Requirements")
    elif block == "sdlc":
        _sdlc_scope_slide(prs, report, workspace, slide_no, title or "SCOPE OF WORK | SDLC Phases")
    elif block == "delivery_effort":
        _delivery_effort_slide(prs, workspace, slide_no, title or "PROJECT DELIVERY | Estimated Effort")
    elif block == "gantt":
        _gantt_slide(prs, _gantt_params_from_wbs(workspace), slide_no,
                     title or "PROJECT DELIVERY | Master Plan & Milestones")
    elif block == "pricing":
        _pricing_slide(prs, report, slide_no, title or "PRICING | CAPEX")
    elif block == "milestones":
        _payment_milestones_slide(prs, slide_no, title or "PRICING | Payment Milestones")
    elif block == "team":
        _team_slide(prs, report, slide_no, title or "PROJECT DELIVERY | Team Structure")
    else:  # safety net
        _detail_slide(prs, title, [], slide_no)


def _append_thank_you(prs: Presentation, thank_you_elements: list) -> None:
    """Append the BnK closing slide.

    Uses the template's dedicated "BnK" closing layout (which carries the brand
    design) so a closing slide is always produced — even when the template's last
    slide had no overlay shapes. Any captured shapes are layered on top.
    """
    slide = prs.slides.add_slide(_layout(prs, "BnK", "C2 -  Separator/ Dark", "Blank"))
    for el in thank_you_elements:
        slide.shapes._spTree.insert_element_before(el, "p:extLst")  # noqa: SLF001


def generate_ppt_proposal_file(
    workspace: Path,
    *,
    title: str = "",
    subtitle: str = "",
    brand: str = "",
    include_sections: list[str] | None = None,
) -> tuple[Path, list[str], list[str]]:
    """Return (pptx_path, sections_rendered, unrecognized_section_names)."""
    template = _template_path()
    if not template.exists():
        raise PPTProposalError(f"BnK template not found: {template}")

    # Reuse existing report data assembly and then apply PPT-specific section names.
    report = assemble_report_data(
        workspace,
        title=title,
        subtitle=subtitle,
        brand=brand,
        include_sections=DEFAULT_REPORT_SECTIONS,
    )
    sections, unrecognized = normalize_ppt_sections(include_sections)

    prs = Presentation(str(template))
    # Save thank-you slide XML before clearing to avoid dangling refs after _clear_slides().
    thank_you_idx = len(prs.slides) - 1 if prs.slides else -1
    thank_you_elements = (
        [deepcopy(s.element) for s in prs.slides[thank_you_idx].shapes]
        if thank_you_idx >= 0 else []
    )
    _clear_slides(prs)

    # --- Prefer the reviewed, traceable storyboard (deck_plan.json) when present ---
    # The deck plan is a CSM projection the user approved at the propose_deck_plan
    # gate; its SlideSpec is a superset of the render outline dict, so we consume it
    # directly. Falls back to the inline LLM outline (then the hardcoded path) below.
    outline: list[dict[str, Any]] | None = None
    outline_source = "hardcoded layout"
    deck_revision: int | None = None
    try:
        from deck import load_deck_plan
        plan = load_deck_plan(workspace)
        if plan and plan.slides:
            deck_revision = plan.revision
            outline = [
                s.model_dump() for s in plan.slides
                if (s.section in sections or s.section in ("", "cover"))
            ]
            outline_source = "deck_plan"
    except Exception as exc:  # noqa: BLE001 — a bad plan must not block rendering
        warnings.warn(f"deck_plan render step failed ({exc!r}); falling back.", stacklevel=2)

    # --- Attempt LLM-driven outline (Paper2Any-style) when no deck plan ---
    if not outline:
        try:
            outline = _generate_slide_outline(report, workspace, sections)
            if outline:
                outline_source = "LLM outline"
        except Exception as exc:
            warnings.warn(f"PPT outline LLM step failed ({exc!r}); using hardcoded layout.", stacklevel=2)

    if outline:
        slide_no = 1
        for spec in outline:
            if spec.get("layout") in CLOSING_LAYOUTS:
                continue  # skip — closing is always appended via _append_thank_you below
            _render_slide(prs, spec, report, workspace, slide_no)
            slide_no += 1
        _append_thank_you(prs, thank_you_elements)
        rendered_sections = sections  # LLM covered all requested sections
    else:
        # --- Fallback: hardcoded section-by-section path ---
        slide_no = 1
        sec = 0  # running section-divider counter for roman numerals
        if "cover" in sections:
            _cover_slide(prs, report, slide_no)
            slide_no += 1
        if "executive_summary" in sections:
            sec += 1
            _section_slide(prs, f"{_roman(sec)}. Executive Summary", slide_no)
            slide_no += 1
            _detail_slide(prs, "EXECUTIVE SUMMARY | Overview", report.get("executive_points", []), slide_no)
            slide_no += 1
        if "solution_overview" in sections:
            sec += 1
            _section_slide(prs, f"{_roman(sec)}. Solution Proposal", slide_no)
            slide_no += 1
            _overview_slide(
                prs,
                f"Solution Proposal\n{report['title']}",
                "Proposed Solution | Scope Of Work | Project Delivery",
                slide_no,
            )
            slide_no += 1
            _detail_slide(
                prs,
                "PROPOSED SOLUTION | Overview",
                [report.get("business_value"), report.get("technical_value"), report.get("pattern_rationale")],
                slide_no,
            )
            slide_no += 1
        if "scope" in sections:
            _functional_nfr_slide(prs, report, slide_no, "PROPOSED SOLUTION | Requirements")
            slide_no += 1
            _sdlc_scope_slide(prs, report, workspace, slide_no)
            slide_no += 1
        if "architecture_diagram" in sections:
            _diagram_slide(prs, report, workspace, slide_no)
            slide_no += 1
        if "technical_stack" in sections:
            _tech_stack_table_slide(prs, report, slide_no)
            slide_no += 1
        if "key_decisions" in sections:
            _detail_slide(prs, "PROPOSED SOLUTION | Key Decisions", _as_list(report["blueprint"].get("key_decisions")), slide_no)
            slide_no += 1
        if "delivery_plan" in sections:
            sec += 1
            _section_slide(prs, f"{_roman(sec)}. Project Delivery", slide_no)
            slide_no += 1
            _delivery_effort_slide(prs, workspace, slide_no)
            slide_no += 1
            _team_slide(prs, report, slide_no)
            slide_no += 1
        if "pricing" in sections:
            sec += 1
            _section_slide(prs, f"{_roman(sec)}. Pricing", slide_no)
            slide_no += 1
            _pricing_slide(prs, report, slide_no)
            slide_no += 1
            _payment_milestones_slide(prs, slide_no)
            slide_no += 1
        if "risks" in sections:
            risk_bullets = [
                f"{r.get('type', 'Risk')}: {r.get('detail', '')} Recommendation: {r.get('recommendation', '')}"
                for r in report.get("risks", [])[:7]
                if isinstance(r, dict)
            ]
            _detail_slide(prs, "PROJECT DELIVERY | Risk & Mitigation", risk_bullets, slide_no)
            slide_no += 1
        if "appendix" in sections:
            artifacts_inv = record_artifact_inventory(workspace)
            artifact_bullets = [f"{a['name']}: {a['label']} ({a['bytes']} bytes)" for a in artifacts_inv]
            _detail_slide(prs, "REFERENCE | Generated Artifacts", artifact_bullets, slide_no)
            slide_no += 1
        _append_thank_you(prs, thank_you_elements)
        rendered_sections = sections

    artifacts = record_artifact_inventory(workspace)
    pptx_path = workspace / "out.pptx"
    prs.save(str(pptx_path))
    slide_count = len(prs.slides)
    record_report_step(
        workspace,
        "generate_ppt_proposal",
        summary=(
            f"Generated editable BnK PowerPoint proposal: {slide_count} slides"
            f" ({outline_source})."
        ),
        data={"sections": rendered_sections, "slide_count": slide_count,
              "artifacts": artifacts, "outline_source": outline_source,
              "deck_plan_revision": deck_revision},
    )
    return pptx_path, rendered_sections, unrecognized

