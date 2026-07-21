import base64
import contextvars
import json

import backends
import session_state as server
import tools
import domain.reporting.reporting as reporting
from tools import GATE_TOOL_NAMES


def _use_workspace(monkeypatch, tmp_path) -> None:
    """Bind ``tmp_path`` as the current-thread workspace (§4.10).

    Stage files resolve lazily against backends.current_workspace(), so swapping the
    ContextVar isolates the whole suite; monkeypatch auto-restores it after the test.
    """
    tmp_path.mkdir(parents=True, exist_ok=True)
    monkeypatch.setattr(
        backends, "_current_workspace",
        contextvars.ContextVar("current_workspace", default=tmp_path),
    )


def _fake_pdf_renderer(html: str, pdf_path) -> None:
    pdf_path.write_bytes(b"%PDF-1.4\n%fake report\n")


def _write_report_inputs(tmp_path) -> None:
    (tmp_path / "architecture_analysis.json").write_text(json.dumps({
        "application_type": "web_application",
        "scale_level": "large",
        "security_level": "high",
        "provider_preference": "aws",
        "detected_capabilities": ["database", "security"],
        "constraints": ["production_focused"],
        "concerns": ["Include observability and security boundaries."],
    }), encoding="utf-8")
    (tmp_path / "diagram_brief.json").write_text(json.dumps({
        "objective": "Deliver a customer-facing architecture for the platform.",
        "functional_requirements": ["Users access the portal", "Application reads relational data"],
        "non_functional_requirements": ["High availability", "Audit-ready security"],
        "layout_constraints": ["Keep runtime and operations concerns separate"],
        "assumptions": ["Traffic volume will be validated during detailed design"],
    }), encoding="utf-8")
    (tmp_path / "tech_stack.json").write_text(json.dumps({
        "frontend": {"choice": "React", "rationale": "Fits SPA delivery", "alternatives": ["Vue"]},
        "database": {"choice": "PostgreSQL", "rationale": "Relational consistency", "alternatives": ["MySQL"]},
    }), encoding="utf-8")
    (tmp_path / "blueprint.json").write_text(json.dumps({
        "slide_title": "Customer Platform Architecture",
        "brand": "Acme",
        "pattern": "three_tier",
        "pattern_rationale": "A three-tier pattern separates access, application, and data layers.",
        "key_decisions": ["Use managed database for operational resilience."],
        "clusters": [{"id": "app", "label": "Application Layer", "tier": "backend"}],
        "nodes": [{"id": "api", "label": "API Service", "tech": "FastAPI", "cluster": "app"}],
        "edges": [{"from": "portal", "to": "api", "label": "HTTPS", "protocol": "HTTP"}],
    }), encoding="utf-8")
    reporting.record_report_step(tmp_path, "test_step", summary="Evidence captured.")


def test_artifacts_includes_pdf_base64(monkeypatch, tmp_path):
    pdf_bytes = b"%PDF-1.4\n%test\n"
    (tmp_path / "out.pdf").write_bytes(pdf_bytes)

    artifacts = server._artifacts(tmp_path)

    assert artifacts["pdf_base64"] == base64.b64encode(pdf_bytes).decode("ascii")

def test_artifacts_includes_pptx_base64(monkeypatch, tmp_path):
    pptx_bytes = b"PK\x03\x04fake pptx"
    (tmp_path / "out.pptx").write_bytes(pptx_bytes)

    artifacts = server._artifacts(tmp_path)

    assert artifacts["pptx_base64"] == base64.b64encode(pptx_bytes).decode("ascii")


def test_generate_pdf_report_maps_to_hitl_card():
    card, step, delta = server._card_for(
        {
            "action_requests": [
                {
                    "name": "generate_pdf_report",
                    "args": {
                        "title": "Architecture Blueprint",
                        "subtitle": "Solution Report",
                        "brand": "Acme",
                        "include_sections": ["cover", "diagram"],
                    },
                }
            ]
        },
        summary="",
    )

    assert "generate_pdf_report" in GATE_TOOL_NAMES
    assert step == "awaiting_pdf_report"
    assert delta == {}
    assert card == {
        "type": "pdf_report_approval",
        "question": "Generate the PDF report with these settings?",
        "title": "Architecture Blueprint",
        "subtitle": "Solution Report",
        "brand": "Acme",
        "include_sections": ["cover", "diagram"],
        "missing_sections": [
            "executive_summary",
            "requirements_analysis",
            "traceability",
            "solution",
            "techstack",
            "architecture_analysis",
            "well_architected",
            "step_results",
            "risks",
        ],
    }

def test_generate_ppt_proposal_maps_to_hitl_card():
    card, step, delta = server._card_for(
        {
            "action_requests": [
                {
                    "name": "generate_ppt_proposal",
                    "args": {
                        "title": "Architecture Proposal",
                        "subtitle": "Solution Deck",
                        "brand": "Acme",
                        "include_sections": ["cover", "architecture_diagram"],
                    },
                }
            ]
        },
        summary="",
    )

    assert "generate_ppt_proposal" in GATE_TOOL_NAMES
    assert step == "awaiting_ppt_proposal"
    assert delta == {}
    assert card["type"] == "ppt_proposal_approval"
    assert card["title"] == "Architecture Proposal"
    assert card["include_sections"] == ["cover", "architecture_diagram"]
    assert "technical_stack" in card["missing_sections"]


def test_tech_stack_gate_persists_pending_draft(monkeypatch, tmp_path):
    _use_workspace(monkeypatch, tmp_path)
    args = {
        "tech_stack": [{"layer": "compute", "choice": "GKE"}],
        "assumptions": {"compliance": ["ISO 27001"]},
        "estimated_total_monthly_cost_usd": {"min_usd": 1000, "max_usd": 2000},
    }

    card, step, delta = server._card_for(
        {"action_requests": [{"name": "propose_tech_stack", "args": args}]},
        summary="",
    )

    assert step == "awaiting_techstack"
    assert card["type"] == "techstack_approval"
    assert delta["tech_stack_draft"] == args
    assert json.loads((tmp_path / "pending_gate.json").read_text(encoding="utf-8"))["tool"] == "propose_tech_stack"
    assert json.loads((tmp_path / "tech_stack_draft.json").read_text(encoding="utf-8")) == args


def test_blueprint_gate_persists_pending_draft(monkeypatch, tmp_path):
    _use_workspace(monkeypatch, tmp_path)
    blueprint = {
        "pattern": "hybrid",
        "nodes": [{"id": "api", "label": "API"}],
        "edges": [{"from": "api", "to": "db"}],
    }

    card, step, delta = server._card_for(
        {"action_requests": [{"name": "propose_blueprint", "args": {"blueprint": blueprint}}]},
        summary="",
    )

    assert step == "awaiting_blueprint"
    assert card["type"] == "blueprint_approval"
    assert delta["blueprint_draft"] == blueprint
    assert json.loads((tmp_path / "pending_gate.json").read_text(encoding="utf-8"))["tool"] == "propose_blueprint"
    assert json.loads((tmp_path / "blueprint_draft.json").read_text(encoding="utf-8")) == blueprint


def test_last_tool_msg_only_when_latest_message_is_tool():
    history_then_user = [
        {"role": "user", "content": "make a diagram"},
        {"role": "tool", "content": '{"approved": true}'},
        {"role": "user", "content": "i want to generate pdf document"},
    ]
    resume = [
        {"role": "user", "content": "make a diagram"},
        {"role": "tool", "content": '{"approved": true}'},
    ]

    assert server._last_tool_msg(history_then_user) is None
    assert server._last_tool_msg(resume) == resume[-1]


def test_pdf_followup_detection():
    assert server._is_pdf_followup("i want to generate pdf document")
    assert server._is_pdf_followup("tạo báo cáo PDF giúp tôi")
    assert not server._is_pdf_followup("please add redis to the diagram")

def test_pdf_followup_does_not_false_positive_on_substrings():
    """Regression test: naive substring matching used to match "doc" inside
    "docker" and "report" inside "reporting", wrongly treating an unrelated
    design request as a PDF follow-up and preserving stale stage artifacts."""
    assert not server._is_pdf_followup("dùng Docker và Kubernetes cho kiến trúc")
    assert not server._is_pdf_followup("add a docker container for the backend")
    assert not server._is_pdf_followup("thêm reporting service vào kiến trúc")

def test_ppt_followup_detection():
    assert server._is_ppt_followup("tạo PPT proposal theo template BnK")
    assert server._is_ppt_followup("make a PowerPoint slide deck")
    assert not server._is_ppt_followup("please add redis to the diagram")


def test_wbs_followup_detection():
    """Regression test: re-export/re-send asks like "xuất lại file WBS" must be
    recognized as a WBS follow-up so chat.py's preserve_artifacts logic skips
    clear_stage_markers() instead of wiping the already-approved wbs.json."""
    assert server._is_wbs_followup("xuất lại file WBS")
    assert server._is_wbs_followup("re-export the WBS excel file")
    assert server._is_wbs_followup("gửi WBS cho khách")
    assert not server._is_wbs_followup("please add redis to the diagram")


def test_wbs_followup_detects_first_time_creation():
    """First-time WBS asks (no wbs.json yet) must ALSO match, else chat.py wipes the
    brief/tech_stack/blueprint that load_solution_context reads and WBS never starts."""
    assert server._is_wbs_followup("tạo WBS cho dự án này")
    assert server._is_wbs_followup("lập kế hoạch công việc giúp tôi")
    assert server._is_wbs_followup("ước lượng effort cho hệ thống")
    assert server._is_wbs_followup("estimate the work breakdown")
    assert server._is_wbs_followup("build the WBS and estimate man-days")
    # unrelated design edits must still NOT match (no false preserve)
    assert not server._is_wbs_followup("please add redis to the diagram")
    assert not server._is_wbs_followup("thêm reporting service vào kiến trúc")


def test_email_followup_detection():
    assert server._is_email_followup("send email now")
    assert server._is_email_followup("gửi qua email cho khách")
    assert server._is_email_followup("gui email cho khach hang")
    assert server._is_email_followup("send this to the client please")
    assert not server._is_email_followup("please add redis to the diagram")


def test_email_followup_does_not_false_positive_on_bare_email_mentions():
    """Regression test: bare "email"/"mail" are common nouns in ordinary design
    requests ("add email verification", "email field on the login form") that
    have nothing to do with sending mail. Matching them would wrongly redirect
    an unrelated design edit into a send_email instruction."""
    assert not server._is_email_followup("add email verification to the signup flow")
    assert not server._is_email_followup("thêm trường email vào form đăng nhập")
    assert not server._is_email_followup("add a mail queue service to the architecture")
    assert not server._is_email_followup("integrate with gmail api for notifications")


def test_email_followup_wins_over_wbs_and_pdf_phrasing():
    """Regression test: "gửi file WBS qua email" / "send the report by email"
    contain "wbs"/"report" and would otherwise be mis-classified by
    _is_wbs_followup/_is_pdf_followup as a re-export/regenerate request. The
    email detector must ALSO fire so chat.py's preserve_email_artifacts branch
    (checked first) wins and the deliverable is sent, not regenerated."""
    msg1 = "gửi file WBS qua email cho khách"
    assert server._is_wbs_followup(msg1)
    assert server._is_email_followup(msg1)

    msg2 = "send the report by email"
    assert server._is_pdf_followup(msg2)
    assert server._is_email_followup(msg2)


def test_wbs_gate_card_coerces_stringified_phases_and_effort():
    """A model that emits phases / effort_by_module / effort_by_role as JSON strings must
    still yield real list/dict card fields, else the frontend Array.isArray/entries guards
    drop them and the WBS approval card renders empty."""
    skeleton_card, _, _ = server._card_for(
        {"action_requests": [{"name": "propose_wbs_skeleton", "args": {
            "phases": '[{"code":"I","name":"SETUP","modules":[{"code":"I.A","name":"Design"}]}]',
        }}]},
        "",
    )
    assert isinstance(skeleton_card["phases"], list)
    assert skeleton_card["phases"][0]["code"] == "I"

    plan_card, _, _ = server._card_for(
        {"action_requests": [{"name": "propose_wbs", "args": {
            "effort_by_module": '[{"code":"II.A","name":"Web","total_md":12}]',
            "effort_by_role": '{"BE": 10, "FE_Mobile": 8}',
        }}]},
        "",
    )
    assert isinstance(plan_card["effort_by_module"], list)
    assert plan_card["effort_by_module"][0]["total_md"] == 12
    assert plan_card["effort_by_role"] == {"BE": 10, "FE_Mobile": 8}


def test_wbs_preserve_first_time_keeps_upstream_artifacts():
    """A first-time WBS request with an upstream solution present preserves artifacts
    (preserve=True) but is NOT already_planned — so chat.py runs the normal wbs_planner
    delegation instead of the re-export shortcut, and never calls clear_stage_markers()."""
    preserve, already = server._wbs_preserve(
        "tạo WBS cho dự án", solution_exists=True, wbs_exists=False, attached=False
    )
    assert preserve is True
    assert already is False


def test_wbs_preserve_reexport_when_wbs_exists():
    preserve, already = server._wbs_preserve(
        "xuất lại WBS", solution_exists=True, wbs_exists=True, attached=False
    )
    assert preserve is True
    assert already is True


def test_wbs_plan_ready_requires_items_and_rollup(tmp_path):
    from routers.chat import _wbs_plan_ready

    (tmp_path / "wbs.json").write_text(
        json.dumps({"items": [], "effort_totals": {"total_mandays": 120}}),
        encoding="utf-8",
    )
    assert _wbs_plan_ready(tmp_path) is False

    (tmp_path / "wbs.json").write_text(
        json.dumps({"items": [{"id": "1.1"}], "effort_totals": {"total_mandays": 0}}),
        encoding="utf-8",
    )
    assert _wbs_plan_ready(tmp_path) is False

    (tmp_path / "wbs.json").write_text(
        json.dumps({"items": [{"id": "1.1"}], "effort_totals": {"total_mandays": 12}}),
        encoding="utf-8",
    )
    assert _wbs_plan_ready(tmp_path) is True


def test_wbs_solution_context_exists_for_direct_render_artifacts(tmp_path):
    from routers.chat import _wbs_solution_context_exists

    assert _wbs_solution_context_exists(tmp_path) is False

    (tmp_path / "requirements.md").write_text("Project notes", encoding="utf-8")
    assert _wbs_solution_context_exists(tmp_path) is True

    (tmp_path / "requirements.md").unlink()
    (tmp_path / "out.png").write_bytes(b"render")
    assert _wbs_solution_context_exists(tmp_path) is True


def test_clear_stage_markers_can_preserve_existing_wbs(monkeypatch, tmp_path):
    _use_workspace(monkeypatch, tmp_path)
    (tmp_path / "blueprint.json").write_text("{}", encoding="utf-8")
    (tmp_path / "wbs_skeleton.json").write_text("{}", encoding="utf-8")
    (tmp_path / "wbs.json").write_text(
        json.dumps({"items": [{"id": "1"}], "effort_totals": {"total_mandays": 12}}),
        encoding="utf-8",
    )
    (tmp_path / "wbs_filled.xlsx").write_bytes(b"xlsx")

    tools.clear_stage_markers(preserve_wbs=True)

    assert not (tmp_path / "blueprint.json").exists()
    assert (tmp_path / "wbs_skeleton.json").exists()
    assert (tmp_path / "wbs.json").exists()
    assert (tmp_path / "wbs_filled.xlsx").exists()


def test_wbs_preserve_no_solution_or_attachment_does_not_preserve():
    # No upstream solution yet -> nothing to preserve (genuine fresh run).
    assert server._wbs_preserve(
        "tạo WBS", solution_exists=False, wbs_exists=False, attached=False
    ) == (False, False)
    # A freshly attached document is new-project intake, never a WBS follow-up.
    assert server._wbs_preserve(
        "tạo WBS", solution_exists=True, wbs_exists=True, attached=True
    ) == (False, False)
    # Non-WBS message never preserves via this path.
    assert server._wbs_preserve(
        "add redis to the diagram", solution_exists=True, wbs_exists=True, attached=False
    ) == (False, False)


def test_wbs_load_solution_context_falls_back_to_requirements(monkeypatch, tmp_path):
    _use_workspace(monkeypatch, tmp_path)
    (tmp_path / "requirements.md").write_text(
        "Project CLARA\n\nBrowser copilot with Planner, Navigator and Validator agents.",
        encoding="utf-8",
    )
    (tmp_path / "out.png").write_bytes(b"png")

    from domain.wbs.wbs_tools import load_solution_context

    digest = json.loads(load_solution_context.func())

    assert digest["source"] == "workspace_fallback"
    assert digest["objective"] == "Project CLARA"
    assert "Browser copilot" in digest["requirements_excerpt"]
    assert "out.png" in digest["available_artifacts"]
    assert "effort_norms" in digest


def test_generate_pdf_report_writes_pdf(monkeypatch, tmp_path):
    from PIL import Image

    _use_workspace(monkeypatch, tmp_path)
    monkeypatch.setattr(reporting, "render_pdf_from_html", _fake_pdf_renderer)
    _write_report_inputs(tmp_path)
    Image.new("RGB", (120, 80), "white").save(tmp_path / "out.png")

    result = tools.generate_pdf_report.func(include_sections=["diagram"])

    assert "Wrote" in result
    assert (tmp_path / "out.report.html").exists()
    assert (tmp_path / "out.pdf").exists()
    assert (tmp_path / "out.pdf").read_bytes().startswith(b"%PDF")

def test_generate_ppt_proposal_writes_openable_pptx(monkeypatch, tmp_path):
    from PIL import Image
    from pptx import Presentation

    _use_workspace(monkeypatch, tmp_path)
    _write_report_inputs(tmp_path)
    Image.new("RGB", (1280, 720), "white").save(tmp_path / "out.png")

    result = tools.generate_ppt_proposal.func(include_sections=["cover", "architecture_diagram"])

    assert "Wrote" in result
    pptx_path = tmp_path / "out.pptx"
    assert pptx_path.exists()
    prs = Presentation(str(pptx_path))
    assert len(prs.slides) >= 2


def test_generate_ppt_proposal_renders_brand_tables(monkeypatch, tmp_path):
    """The fallback path (no LLM) must emit native, brand-styled tables, not just bullets."""
    from PIL import Image
    from pptx import Presentation

    _use_workspace(monkeypatch, tmp_path)
    _write_report_inputs(tmp_path)
    Image.new("RGB", (1280, 720), "white").save(tmp_path / "out.png")

    result = tools.generate_ppt_proposal.func(
        include_sections=["cover", "technical_stack", "scope", "pricing"]
    )

    assert "Wrote" in result
    prs = Presentation(str(tmp_path / "out.pptx"))
    # At least one slide must contain a real table (graphic frame), proving table builders ran.
    tables = [
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_table
    ]
    assert tables, "expected at least one native table in the proposal"


def test_report_data_uses_step_results_and_section_aliases(tmp_path):
    from PIL import Image

    _write_report_inputs(tmp_path)
    Image.new("RGB", (120, 80), "white").save(tmp_path / "out.png")

    data = reporting.assemble_report_data(tmp_path, include_sections=["cover", "blueprint", "diagram"])

    assert data["sections"] == ["cover", "architecture_analysis", "diagram"]
    assert data["title"] == "Customer Platform Architecture"
    assert data["analysis"]["security_level"] == "high"
    assert data["evidence_steps"][0]["step"] == "test_step"
    assert data["traceability"]


def test_report_data_falls_back_to_rendered_workspace_artifacts(tmp_path):
    from PIL import Image

    (tmp_path / "requirements.md").write_text(
        """
        <untrusted_document>
        1.1 This document outlines the requirements and scope for the design,
        development, integration, and implementation of new modules for the HIVE system.
        2.1.1 The web-based portal should provide summary views and data extraction.
        2.1.1 The mobile application shall allow simplified functionality for field users.
        4.1 The system shall include audit-ready security, MFA, and operational logging.
        </untrusted_document>
        """,
        encoding="utf-8",
    )
    (tmp_path / "layout_plan.json").write_text(
        json.dumps({"notes": ["bundled repetitive hub edges"]}),
        encoding="utf-8",
    )
    (tmp_path / "out.native_stats.json").write_text(
        json.dumps({"nodes": 42, "edges": 40, "bundled_edges": 26}),
        encoding="utf-8",
    )
    (tmp_path / "out.drawio").write_text(
        """
        <mxfile compressed="false">
          <diagram name="CAG HIVE AWS Architecture">
            <mxGraphModel><root>
              <mxCell id="0" />
              <mxCell id="1" parent="0" />
              <mxCell id="__title" value="HIVE user journeys, integrations, data and operations" vertex="1" parent="1" />
              <mxCell id="web_users" value="Web Users&lt;br&gt;Browser / HTTPS" vertex="1" parent="1" />
              <mxCell id="waf" value="AWS WAF&lt;br&gt;AWS WAF managed rules" vertex="1" parent="1" />
              <mxCell id="api_gateway" value="API Gateway&lt;br&gt;Amazon API Gateway private APIs" vertex="1" parent="1" />
              <mxCell id="odc_runtime" value="OutSystems ODC Runtime&lt;br&gt;Production capacity" vertex="1" parent="1" />
              <mxCell id="cloudwatch" value="Operations Monitoring&lt;br&gt;Amazon CloudWatch" vertex="1" parent="1" />
              <mxCell id="e1" value="HTTPS request" edge="1" source="web_users" target="waf" parent="1" />
              <mxCell id="e2" value="private API" edge="1" source="waf" target="api_gateway" parent="1" />
            </root></mxGraphModel>
          </diagram>
        </mxfile>
        """,
        encoding="utf-8",
    )
    Image.new("RGB", (120, 80), "white").save(tmp_path / "out.png")

    data = reporting.assemble_report_data(tmp_path)
    html = reporting.render_report_html(data)

    assert data["title"] == "CAG HIVE AWS Architecture"
    assert data["node_count"] == 42
    assert data["edge_count"] == 40
    assert data["brief"]["functional_requirements"]
    assert data["analysis"]["provider_preference"] == "aws"
    assert any(item["choice"] == "AWS WAF" for item in data["tech_items"])
    assert data["blueprint"]["nodes"][0]["label"] == "Web Users"
    assert "No functional requirements were captured" not in html
    assert "No technology stack was recorded" not in html

def test_report_template_escapes_user_text(tmp_path):
    from PIL import Image

    _write_report_inputs(tmp_path)
    brief = json.loads((tmp_path / "diagram_brief.json").read_text(encoding="utf-8"))
    brief["objective"] = "<script>alert('x')</script>"
    (tmp_path / "diagram_brief.json").write_text(json.dumps(brief), encoding="utf-8")
    Image.new("RGB", (120, 80), "white").save(tmp_path / "out.png")

    html = reporting.render_report_html(reporting.assemble_report_data(tmp_path))

    assert "<script>" not in html
    assert "&lt;script&gt;" in html
    assert "Executive Summary" in html
    assert "Step Results and Quality Gates" in html


def test_generate_pdf_report_missing_playwright_is_actionable(monkeypatch, tmp_path):
    from PIL import Image

    _use_workspace(monkeypatch, tmp_path)
    _write_report_inputs(tmp_path)
    Image.new("RGB", (120, 80), "white").save(tmp_path / "out.png")

    def fail_renderer(html: str, pdf_path) -> None:
        raise reporting.ReportRenderError("Playwright Chromium is not available.")

    monkeypatch.setattr(reporting, "render_pdf_from_html", fail_renderer)

    result = tools.generate_pdf_report.func(include_sections=["diagram"])

    assert "PDF report generation failed" in result
    assert "Playwright Chromium" in result
